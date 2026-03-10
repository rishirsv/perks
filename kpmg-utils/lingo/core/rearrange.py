"""Slide duplication utilities for KPMG-PPTX.

The duplicate_slide function performs a deep copy of a slide while
preserving image and other relationships. It is a focused port of the
proven logic used in the document-skills toolkit, trimmed to the
minimal surface we need for template cloning.
"""

from copy import deepcopy
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
REL_PREFIX = f"{{{REL_NS}}}"


def _clone_relationships(src_part, dest_part):
    """Copy relationships from ``src_part`` to ``dest_part``.

    Returns a mapping of old relationship IDs to the newly created IDs so
    we can update embedded references (e.g., ``r:embed`` on blips).

    Some python-pptx versions implement ``_Relationships.values()`` via
    ``[self[k] for k in self]`` while ``__iter__`` yields relationship
    objects, not rId strings, which triggers a KeyError lookup. To avoid
    that, prefer the underlying ``_rels`` dict when available.
    """

    rel_id_map: dict[str, str] = {}

    rels_collection = src_part.rels
    raw_rels = getattr(rels_collection, "_rels", None)

    # Prefer the plain dict to avoid Mapping.values() relying on buggy __iter__
    if raw_rels is not None:
        rels_iter = list(raw_rels.values())
    else:
        # Fallback for future python-pptx versions
        rels_iter = list(rels_collection)

    # Force resolution of all relationships up front to avoid lazy-load bugs
    for rel in rels_iter:
        # Skip notesSlide/slide master rels; python-pptx manages them
        if rel.reltype in {RT.SLIDE_MASTER, RT.NOTES_SLIDE, RT.SLIDE_LAYOUT}:
            continue

        old_rid = rel.rId
        new_rid = dest_part.rels._add_relationship(rel.reltype, rel._target, rel.is_external)
        rel_id_map[old_rid] = new_rid

    return rel_id_map


def _remap_blip_rels(element, rel_id_map):
    """Update all ``a:blip`` embed references in ``element`` using ``rel_id_map``."""

    nsmap = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }

    for blip in element.iterfind('.//a:blip', namespaces=nsmap):
        embed = blip.get(f"{{{nsmap['r']}}}embed")
        if embed and embed in rel_id_map:
            blip.set(f"{{{nsmap['r']}}}embed", rel_id_map[embed])


def _remap_relationship_ids(element, rel_id_map):
    """Update any r:* attributes in ``element`` using ``rel_id_map``.

    This catches embeds, hyperlinks, OLE objects, and any other relationship-
    backed references that appear on shapes. Without this, duplicated slides
    can retain stale relationship IDs when the newly created relationships use
    different identifiers (common in real-world templates with non-contiguous
    rIds), which leads to missing-relationship errors on save.
    """

    if not rel_id_map:
        return

    for el in element.iter():
        for attr, val in list(el.attrib.items()):
            if attr.startswith(REL_PREFIX) and val in rel_id_map:
                el.attrib[attr] = rel_id_map[val]


def duplicate_slide(presentation, slide_index):
    """Duplicate a slide in ``presentation`` and return the new slide.

    This deep-copies the slide XML, preserves ordering, and recreates the
    embedded image/chart relationships so copied shapes continue to
    display correctly.

    Args:
        presentation: python-pptx ``Presentation`` instance.
        slide_index: int, zero-based index of the slide to duplicate.

    Returns:
        The newly created slide object (appended to the presentation).
    """

    src = presentation.slides[slide_index]
    dest = presentation.slides.add_slide(src.slide_layout)

    # Remove all auto-generated placeholders from the new slide
    # This prevents "double shapes" (one from layout, one copied from source)
    for shape in list(dest.shapes):
        dest.shapes._spTree.remove(shape.element)

    # Copy relationships first so we can remap embeds when cloning shapes
    rel_id_map = _clone_relationships(src.part, dest.part)

    # Deep-copy every shape (maintain original ordering)
    for shape in src.shapes:
        new_el = deepcopy(shape.element)
        _remap_relationship_ids(new_el, rel_id_map)
        dest.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    return dest


def delete_slide(presentation, index: int) -> None:
    """Delete slide at given index from presentation.

    Args:
        presentation: python-pptx Presentation instance
        index: Zero-based slide index to delete

    Raises:
        IndexError: If index is out of range
    """
    if index < 0 or index >= len(presentation.slides):
        raise IndexError(f"Slide index {index} out of range (0-{len(presentation.slides)-1})")

    rId = presentation.slides._sldIdLst[index].rId
    presentation.part.drop_rel(rId)
    del presentation.slides._sldIdLst[index]


__all__ = ["duplicate_slide", "delete_slide"]
