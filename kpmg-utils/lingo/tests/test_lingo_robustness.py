"""Robustness tests for Lingo glossary generator.

These tests verify basic functionality and error handling. Note that with
dynamic pagination, slide counts and distribution depend on content length
estimation, not fixed row counts.
"""

import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from core.glossary_generator import generate_glossary


class LingoRobustnessTestCase(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls.template = Path("templates/glossary_template.pptx")
        prs = Presentation(cls.template)
        tables = [shape.table for shape in prs.slides[0].shapes if shape.has_table]
        cls.left_cap = len(tables[0].rows) - 1
        cls.right_cap = len(tables[1].rows) - 1
        # Dynamic pagination uses min of left/right as rows_per_column
        cls.rows_per_column = min(cls.left_cap, cls.right_cap)

    def _make_terms(self, count):
        return [
            {"term": f"T{idx:03d}", "definition": f"Definition {idx}"}
            for idx in range(count)
        ]

    def _count_filled_rows(self, table):
        return sum(1 for row_idx in range(1, len(table.rows)) if table.cell(row_idx, 0).text.strip())

    def _total_filled_terms(self, prs):
        """Count total filled terms across all slides."""
        total = 0
        for slide in prs.slides:
            tables = [shape.table for shape in slide.shapes if shape.has_table]
            for table in tables[:2]:
                total += self._count_filled_rows(table)
        return total

    def _generate_and_load(self, count):
        with tempfile.TemporaryDirectory() as tmpdir:
            output = Path(tmpdir) / "out.pptx"
            generate_glossary(str(self.template), self._make_terms(count), str(output))
            return Presentation(output)

    def test_all_terms_preserved(self):
        """All input terms should appear in output, regardless of pagination."""
        for count in [10, 30, 34, 35, 50, 70]:
            with self.subTest(term_count=count):
                prs = self._generate_and_load(count)
                total_filled = self._total_filled_terms(prs)
                self.assertEqual(total_filled, count)

    def test_short_terms_fit_efficiently(self):
        """Short terms should fit to the template row capacity."""
        capacity = self.left_cap + self.right_cap

        # Below capacity should fit on 1 slide
        prs = self._generate_and_load(max(1, capacity - 1))
        self.assertEqual(len(prs.slides), 1)

        # Exact capacity should fit on 1 slide
        prs = self._generate_and_load(capacity)
        self.assertEqual(len(prs.slides), 1)

        # One over capacity should require multiple slides
        prs = self._generate_and_load(capacity + 1)
        self.assertGreaterEqual(len(prs.slides), 2)

    def test_each_slide_has_two_tables(self):
        """Each generated slide should have exactly 2 tables."""
        prs = self._generate_and_load(50)
        for slide_idx, slide in enumerate(prs.slides):
            with self.subTest(slide=slide_idx):
                tables = [shape.table for shape in slide.shapes if shape.has_table]
                self.assertEqual(len(tables), 2)

    def test_invalid_template_raises_value_error(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            bad_template = Path(tmpdir) / "bad_template.pptx"
            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[6])  # blank, no tables
            prs.save(bad_template)

            with self.assertRaisesRegex(ValueError, "Template must contain at least 2 tables"):
                generate_glossary(str(bad_template), self._make_terms(5), str(Path(tmpdir) / "out.pptx"))


if __name__ == "__main__":
    unittest.main()
