"""
Value Driver Diagram Generator for FDD Research

Generates PowerPoint value driver tree diagrams for financial due diligence.
Designed to run in ChatGPT's Code Interpreter environment.

Usage:
    from value_driver_diagram import create_value_driver_pptx

    pptx_path = create_value_driver_pptx(
        industry="Technology / SaaS",
        company_name="Acme Software",
        metrics={"ARR": "$42M", "Growth": "+28%", "Gross Margin": "78%"}
    )
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import re
from typing import Dict, List, Optional, Tuple
from dataclasses import dataclass


# ============================================================================
# Color Palette
# ============================================================================

COLORS = {
    "primary": RGBColor(0x1a, 0x36, 0x5d),      # Dark blue
    "secondary": RGBColor(0x2e, 0x5a, 0x87),    # Medium blue
    "accent": RGBColor(0x4a, 0x90, 0xd9),       # Light blue
    "highlight": RGBColor(0xff, 0x9f, 0x1c),    # Orange (for key drivers)
    "risk": RGBColor(0xdc, 0x35, 0x45),         # Red (for risks)
    "text_dark": RGBColor(0x2d, 0x3a, 0x4a),    # Dark text
    "text_light": RGBColor(0xff, 0xff, 0xff),   # White text
    "bg_light": RGBColor(0xf8, 0xf9, 0xfa),     # Light background
    "border": RGBColor(0xde, 0xe2, 0xe6),       # Border gray
}


# ============================================================================
# Industry Tree Definitions
# ============================================================================

@dataclass
class TreeNode:
    """Represents a node in the value driver tree."""
    label: str
    children: List['TreeNode']
    is_key_driver: bool = False
    is_risk: bool = False
    metric: Optional[str] = None


INDUSTRY_TREES = {
    "Technology / SaaS": {
        "ebitda": TreeNode("EBITDA", [
            TreeNode("Net Revenue", [
                TreeNode("Subscription Revenue", [
                    TreeNode("Starting ARR", [], is_key_driver=True),
                    TreeNode("+ New ARR", []),
                    TreeNode("+ Expansion ARR", [], is_key_driver=True),
                    TreeNode("- Contraction", []),
                    TreeNode("- Churn", [], is_risk=True),
                ]),
                TreeNode("Usage Revenue", []),
                TreeNode("Services Revenue", []),
                TreeNode("Gross-to-Net", [], is_risk=True),
            ]),
            TreeNode("COGS", [
                TreeNode("Hosting/Cloud", []),
                TreeNode("Support/Success", []),
                TreeNode("Third-party costs", []),
            ]),
            TreeNode("OpEx", [
                TreeNode("Sales & Marketing", []),
                TreeNode("R&D", []),
                TreeNode("G&A", []),
            ]),
        ]),
        "working_capital": TreeNode("Working Capital", [
            TreeNode("Receivables (A/R)", [
                TreeNode("DSO", []),
                TreeNode("Customer mix", []),
                TreeNode("Billing frequency", []),
            ]),
            TreeNode("Deferred Revenue", [
                TreeNode("Annual vs Monthly", [], is_key_driver=True),
                TreeNode("Multi-year deals", []),
            ], is_key_driver=True),
            TreeNode("Payables (A/P)", [
                TreeNode("Cloud vendor terms", []),
            ]),
            TreeNode("Accruals", [
                TreeNode("Commissions", []),
                TreeNode("Bonus", []),
            ]),
        ]),
        "capex": TreeNode("Capex", [
            TreeNode("Growth Capex", [
                TreeNode("Capitalized Software", [], is_key_driver=True),
                TreeNode("Infrastructure", []),
            ]),
            TreeNode("Maintenance Capex", [
                TreeNode("Hardware", []),
                TreeNode("Office", []),
            ]),
        ]),
    },

    "Banking & Lending": {
        "ebitda": TreeNode("PPNR", [
            TreeNode("Net Interest Income", [
                TreeNode("Interest Income", [
                    TreeNode("Loan Yield", [], is_key_driver=True),
                    TreeNode("Securities Yield", []),
                ]),
                TreeNode("Interest Expense", [
                    TreeNode("Deposit Cost", [], is_key_driver=True),
                    TreeNode("Borrowing Cost", []),
                ]),
            ]),
            TreeNode("Non-Interest Income", [
                TreeNode("Origination Fees", []),
                TreeNode("Servicing Fees", []),
                TreeNode("Gain on Sale", [], is_risk=True),
            ]),
            TreeNode("Operating Expenses", [
                TreeNode("Compensation", []),
                TreeNode("Technology", []),
                TreeNode("Occupancy", []),
            ]),
            TreeNode("Provision Expense", [], is_risk=True),
        ]),
        "working_capital": TreeNode("Balance Sheet Drivers", [
            TreeNode("Loan Portfolio", [
                TreeNode("Origination Volume", [], is_key_driver=True),
                TreeNode("Payoffs/Paydowns", []),
                TreeNode("Charge-offs", [], is_risk=True),
            ]),
            TreeNode("Deposits", [
                TreeNode("Core Deposits", [], is_key_driver=True),
                TreeNode("Brokered Deposits", [], is_risk=True),
            ]),
            TreeNode("Warehouse Lines", []),
        ]),
        "capex": TreeNode("Capital & Reserves", [
            TreeNode("ACL/CECL Reserve", [], is_risk=True),
            TreeNode("Regulatory Capital", []),
            TreeNode("Technology Investment", []),
        ]),
    },

    "Transportation & Logistics": {
        "ebitda": TreeNode("EBITDA", [
            TreeNode("Revenue", [
                TreeNode("Trucking", [
                    TreeNode("# Loads", [], is_key_driver=True),
                    TreeNode("x Revenue/Load", [], is_key_driver=True),
                ]),
                TreeNode("Brokerage", [
                    TreeNode("Gross Revenue", []),
                    TreeNode("- Purchased Transportation", [], is_risk=True),
                ]),
                TreeNode("Warehousing", []),
                TreeNode("Fuel Surcharge", []),
            ]),
            TreeNode("Operating Costs", [
                TreeNode("Driver Wages", [], is_key_driver=True),
                TreeNode("Fuel (net)", []),
                TreeNode("Equipment", []),
                TreeNode("Insurance", [], is_risk=True),
                TreeNode("Maintenance", []),
            ]),
            TreeNode("SG&A", []),
        ]),
        "working_capital": TreeNode("Working Capital", [
            TreeNode("Receivables", [
                TreeNode("DSO by customer", [], is_key_driver=True),
                TreeNode("Broker vs direct", []),
            ]),
            TreeNode("Payables", [
                TreeNode("Carrier payments", []),
                TreeNode("Fuel terms", []),
            ]),
        ]),
        "capex": TreeNode("Capex", [
            TreeNode("Fleet Capex", [
                TreeNode("Tractors", [], is_key_driver=True),
                TreeNode("Trailers", []),
                TreeNode("Lease vs Own", []),
            ]),
            TreeNode("Maintenance Capex", []),
            TreeNode("Technology (TMS)", []),
        ]),
    },

    "Insurance": {
        "ebitda": TreeNode("Underwriting Income", [
            TreeNode("Net Earned Premium", [
                TreeNode("Gross Written Premium", [], is_key_driver=True),
                TreeNode("- Ceded Premium", []),
                TreeNode("Unearned Premium Change", []),
            ]),
            TreeNode("Losses & LAE", [
                TreeNode("Paid Losses", []),
                TreeNode("Reserve Change", [], is_risk=True),
                TreeNode("IBNR", [], is_risk=True),
            ]),
            TreeNode("Underwriting Expense", [
                TreeNode("Commissions", []),
                TreeNode("Operating Expenses", []),
            ]),
            TreeNode("Investment Income", [
                TreeNode("Portfolio Yield", []),
                TreeNode("Realized Gains/Losses", [], is_risk=True),
            ]),
        ]),
        "working_capital": TreeNode("Float & Reserves", [
            TreeNode("Loss Reserves", [
                TreeNode("Case Reserves", []),
                TreeNode("IBNR", [], is_risk=True),
            ]),
            TreeNode("Unearned Premium", [], is_key_driver=True),
            TreeNode("Reinsurance Recoverables", []),
        ]),
        "capex": TreeNode("Capital", [
            TreeNode("Statutory Surplus", [], is_key_driver=True),
            TreeNode("RBC Requirements", []),
            TreeNode("Technology Investment", []),
        ]),
    },

    "Real Estate & Construction": {
        "ebitda": TreeNode("NOI / EBITDA", [
            TreeNode("Revenue", [
                TreeNode("Base Rent", [
                    TreeNode("Occupied SF", [], is_key_driver=True),
                    TreeNode("x Rent/SF", [], is_key_driver=True),
                ]),
                TreeNode("CAM Recoveries", []),
                TreeNode("Other Income", []),
                TreeNode("Construction Revenue", [
                    TreeNode("Backlog", [], is_key_driver=True),
                    TreeNode("x POC %", []),
                ]),
            ]),
            TreeNode("Property Expenses", [
                TreeNode("Property Taxes", []),
                TreeNode("Insurance", []),
                TreeNode("Utilities", []),
                TreeNode("Repairs & Maint", []),
            ]),
            TreeNode("Construction Costs", [
                TreeNode("Direct Costs", []),
                TreeNode("Subcontractor", []),
            ]),
        ]),
        "working_capital": TreeNode("Working Capital", [
            TreeNode("Receivables", [
                TreeNode("Tenant A/R", []),
                TreeNode("Retainage", [], is_key_driver=True),
            ]),
            TreeNode("CIP / WIP", [], is_key_driver=True),
            TreeNode("Billings", [
                TreeNode("Over-billings", [], is_risk=True),
                TreeNode("Under-billings", []),
            ]),
            TreeNode("Payables", [
                TreeNode("Subcontractor A/P", []),
                TreeNode("Retention payable", []),
            ]),
        ]),
        "capex": TreeNode("Capex", [
            TreeNode("Development", [], is_key_driver=True),
            TreeNode("TI/LC", []),
            TreeNode("Maintenance", []),
            TreeNode("Equipment", []),
        ]),
    },

    "Business Services": {
        "ebitda": TreeNode("EBITDA", [
            TreeNode("Revenue", [
                TreeNode("T&M Revenue", [
                    TreeNode("Billable Hours", [], is_key_driver=True),
                    TreeNode("x Bill Rate", [], is_key_driver=True),
                ]),
                TreeNode("Fixed-Fee Revenue", [
                    TreeNode("# Projects", []),
                    TreeNode("x Avg Size", []),
                    TreeNode("x Completion %", []),
                ]),
                TreeNode("Managed Services", []),
                TreeNode("Reimbursables", []),
            ]),
            TreeNode("Cost of Services", [
                TreeNode("Consultant Comp", [], is_key_driver=True),
                TreeNode("Subcontractors", []),
                TreeNode("Travel/Expenses", []),
            ]),
            TreeNode("SG&A", [
                TreeNode("Sales/BD", []),
                TreeNode("Corporate", []),
            ]),
        ]),
        "working_capital": TreeNode("Working Capital", [
            TreeNode("Receivables", [
                TreeNode("DSO", [], is_key_driver=True),
                TreeNode("Unbilled WIP", [], is_risk=True),
            ]),
            TreeNode("Deferred Revenue", []),
            TreeNode("Payables", [
                TreeNode("Subcontractor A/P", []),
            ]),
            TreeNode("Accruals", [
                TreeNode("Bonus accrual", []),
                TreeNode("PTO accrual", []),
            ]),
        ]),
        "capex": TreeNode("Capex", [
            TreeNode("Technology", []),
            TreeNode("Office Build-out", []),
            TreeNode("Training", []),
        ]),
    },

    "Healthcare Services": {
        "ebitda": TreeNode("EBITDA", [
            TreeNode("Net Patient Revenue", [
                TreeNode("Gross Charges", [
                    TreeNode("Visits/Encounters", [], is_key_driver=True),
                    TreeNode("x wRVUs/Visit", []),
                    TreeNode("x Charge/wRVU", []),
                ]),
                TreeNode("Contractual Adjustments", [
                    TreeNode("Medicare rates", []),
                    TreeNode("Medicaid rates", []),
                    TreeNode("Commercial rates", [], is_key_driver=True),
                ]),
                TreeNode("Bad Debt", [], is_risk=True),
            ]),
            TreeNode("Direct Costs", [
                TreeNode("Provider Comp", [], is_key_driver=True),
                TreeNode("Clinical Staff", []),
                TreeNode("Supplies", []),
            ]),
            TreeNode("Facility Costs", [
                TreeNode("Rent", []),
                TreeNode("Utilities", []),
            ]),
            TreeNode("SG&A", []),
        ]),
        "working_capital": TreeNode("Working Capital", [
            TreeNode("Patient A/R", [
                TreeNode("DSO by payor", [], is_key_driver=True),
                TreeNode("Denials", [], is_risk=True),
            ]),
            TreeNode("Payables", []),
            TreeNode("Accrued Comp", []),
        ]),
        "capex": TreeNode("Capex", [
            TreeNode("Medical Equipment", []),
            TreeNode("EHR/Technology", []),
            TreeNode("Facility Expansion", [], is_key_driver=True),
        ]),
    },

    "Retail / Consumer": {
        "ebitda": TreeNode("EBITDA", [
            TreeNode("Revenue", [
                TreeNode("Store Revenue", [
                    TreeNode("# Stores", [], is_key_driver=True),
                    TreeNode("x Sales/Store", [], is_key_driver=True),
                    TreeNode("Comp Sales", [], is_key_driver=True),
                ]),
                TreeNode("E-commerce", [
                    TreeNode("Traffic", []),
                    TreeNode("x Conversion", []),
                    TreeNode("x AOV", []),
                ]),
                TreeNode("Wholesale", []),
            ]),
            TreeNode("COGS", [
                TreeNode("Product Cost", [], is_key_driver=True),
                TreeNode("Freight In", []),
                TreeNode("Shrink", [], is_risk=True),
            ]),
            TreeNode("Store OpEx", [
                TreeNode("Labor", [], is_key_driver=True),
                TreeNode("Occupancy", []),
                TreeNode("Other", []),
            ]),
            TreeNode("Corporate", []),
        ]),
        "working_capital": TreeNode("Working Capital", [
            TreeNode("Inventory", [
                TreeNode("DIO / Turns", [], is_key_driver=True),
                TreeNode("Aging", [], is_risk=True),
            ]),
            TreeNode("Receivables", [
                TreeNode("Wholesale A/R", []),
            ]),
            TreeNode("Payables", [
                TreeNode("Vendor terms", [], is_key_driver=True),
            ]),
        ]),
        "capex": TreeNode("Capex", [
            TreeNode("New Stores", [], is_key_driver=True),
            TreeNode("Remodels", []),
            TreeNode("Technology", []),
            TreeNode("DC Investment", []),
        ]),
    },

    "Industrial Manufacturing": {
        "ebitda": TreeNode("EBITDA", [
            TreeNode("Revenue", [
                TreeNode("Product Revenue", [
                    TreeNode("Units Shipped", [], is_key_driver=True),
                    TreeNode("x ASP", [], is_key_driver=True),
                ]),
                TreeNode("Aftermarket/Parts", []),
                TreeNode("Services", []),
            ]),
            TreeNode("COGS", [
                TreeNode("Materials", [], is_key_driver=True),
                TreeNode("Direct Labor", []),
                TreeNode("Manufacturing OH", []),
                TreeNode("Scrap/Yield Loss", [], is_risk=True),
            ]),
            TreeNode("SG&A", []),
            TreeNode("R&D", []),
        ]),
        "working_capital": TreeNode("Working Capital", [
            TreeNode("Receivables", [
                TreeNode("DSO", []),
            ]),
            TreeNode("Inventory", [
                TreeNode("Raw Materials", []),
                TreeNode("WIP", []),
                TreeNode("Finished Goods", [], is_key_driver=True),
                TreeNode("Obsolescence", [], is_risk=True),
            ]),
            TreeNode("Payables", [
                TreeNode("Supplier terms", [], is_key_driver=True),
            ]),
        ]),
        "capex": TreeNode("Capex", [
            TreeNode("Capacity Expansion", [], is_key_driver=True),
            TreeNode("Automation", []),
            TreeNode("Maintenance", []),
            TreeNode("Tooling", []),
        ]),
    },

    "Asset Management": {
        "ebitda": TreeNode("Revenue & Earnings", [
            TreeNode("Management Fees", [
                TreeNode("AUM", [], is_key_driver=True),
                TreeNode("x Fee Rate", [], is_key_driver=True),
            ]),
            TreeNode("Performance Fees", [
                TreeNode("Carry/Incentive", [], is_risk=True),
                TreeNode("Hurdle Rate", []),
            ]),
            TreeNode("Transaction Fees", []),
            TreeNode("Operating Expenses", [
                TreeNode("Compensation", [], is_key_driver=True),
                TreeNode("Distribution", []),
                TreeNode("Technology", []),
                TreeNode("G&A", []),
            ]),
        ]),
        "working_capital": TreeNode("Working Capital", [
            TreeNode("Fee Receivables", []),
            TreeNode("Accrued Carry", [], is_key_driver=True),
            TreeNode("GP Commitments", []),
            TreeNode("Seed Capital", []),
        ]),
        "capex": TreeNode("Capital", [
            TreeNode("GP Commitments", [], is_key_driver=True),
            TreeNode("Balance Sheet Investments", []),
            TreeNode("Technology Platform", []),
        ]),
    },
}


# ============================================================================
# PowerPoint Generation
# ============================================================================

def add_title_slide(prs: Presentation, company_name: str, industry: str):
    """Add title slide to presentation."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS["primary"]
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(9), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{company_name}"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS["text_light"]
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.6), Inches(9), Inches(0.6)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Value Driver Analysis"
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS["accent"]
    p.alignment = PP_ALIGN.CENTER

    # Industry tag
    industry_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.4), Inches(9), Inches(0.5)
    )
    tf = industry_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Industry: {industry}"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS["text_light"]
    p.alignment = PP_ALIGN.CENTER


def draw_tree_node(
    slide,
    node: TreeNode,
    x: float,
    y: float,
    width: float,
    height: float,
    level: int = 0,
    metrics: Dict[str, str] = None
) -> Tuple[float, float]:
    """Draw a single node and return its center position."""
    metrics = metrics or {}

    # Node styling based on type
    if node.is_key_driver:
        fill_color = COLORS["highlight"]
        text_color = COLORS["text_dark"]
    elif node.is_risk:
        fill_color = COLORS["risk"]
        text_color = COLORS["text_light"]
    elif level == 0:
        fill_color = COLORS["primary"]
        text_color = COLORS["text_light"]
    elif level == 1:
        fill_color = COLORS["secondary"]
        text_color = COLORS["text_light"]
    else:
        fill_color = COLORS["accent"]
        text_color = COLORS["text_dark"]

    # Create shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = COLORS["border"]

    # Add text
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]

    # Check if we have a metric for this node
    label = node.label
    if node.metric:
        label = f"{label}\n{node.metric}"
    elif metrics:
        for key, value in metrics.items():
            if key.lower() in label.lower():
                label = f"{label}\n{value}"
                break

    p.text = label
    p.font.size = Pt(9) if level > 1 else Pt(10)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

    # Vertical alignment
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return (x + width/2, y + height/2)


def draw_connector(slide, start: Tuple[float, float], end: Tuple[float, float]):
    """Draw a connector line between two points."""
    connector = slide.shapes.add_connector(
        1,  # Straight connector
        Inches(start[0]), Inches(start[1]),
        Inches(end[0]), Inches(end[1])
    )
    connector.line.color.rgb = COLORS["border"]
    connector.line.width = Pt(1)


def add_tree_slide(
    prs: Presentation,
    title: str,
    tree: TreeNode,
    metrics: Dict[str, str] = None
):
    """Add a slide with a value driver tree."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS["text_dark"]

    # Draw tree using hierarchical layout
    draw_tree_hierarchical(slide, tree, metrics)

    # Legend
    add_legend(slide)


def draw_tree_hierarchical(
    slide,
    root: TreeNode,
    metrics: Dict[str, str] = None,
    start_x: float = 0.3,
    start_y: float = 0.9,
    total_width: float = 9.4,
    level_height: float = 1.1
):
    """Draw tree in hierarchical layout."""
    metrics = metrics or {}

    # Calculate tree structure
    def count_leaves(node: TreeNode) -> int:
        if not node.children:
            return 1
        return sum(count_leaves(c) for c in node.children)

    def draw_subtree(
        node: TreeNode,
        x: float,
        y: float,
        available_width: float,
        level: int
    ) -> Tuple[float, float]:
        """Recursively draw subtree, return center of root node."""

        node_width = min(1.4, available_width * 0.8)
        node_height = 0.45

        if not node.children:
            # Leaf node
            center_x = x + available_width / 2
            node_x = center_x - node_width / 2
            center = draw_tree_node(
                slide, node, node_x, y, node_width, node_height, level, metrics
            )
            return center

        # Non-leaf: draw children first to get their positions
        child_centers = []
        total_leaves = count_leaves(node)
        current_x = x

        for child in node.children:
            child_leaves = count_leaves(child)
            child_width = (child_leaves / total_leaves) * available_width
            child_center = draw_subtree(
                child, current_x, y + level_height, child_width, level + 1
            )
            child_centers.append(child_center)
            current_x += child_width

        # Position parent centered above children
        if child_centers:
            center_x = (child_centers[0][0] + child_centers[-1][0]) / 2
        else:
            center_x = x + available_width / 2

        node_x = center_x - node_width / 2
        parent_center = draw_tree_node(
            slide, node, node_x, y, node_width, node_height, level, metrics
        )

        # Draw connectors
        for child_center in child_centers:
            # Draw from bottom of parent to top of child
            draw_connector(
                slide,
                (parent_center[0], parent_center[1] + node_height/2),
                (child_center[0], y + level_height)
            )

        return parent_center

    draw_subtree(root, start_x, start_y, total_width, 0)


def add_legend(slide):
    """Add legend to slide."""
    legend_y = 6.8

    # Key Driver legend
    kd_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.5), Inches(legend_y),
        Inches(0.3), Inches(0.2)
    )
    kd_shape.fill.solid()
    kd_shape.fill.fore_color.rgb = COLORS["highlight"]
    kd_shape.line.fill.background()

    kd_label = slide.shapes.add_textbox(
        Inches(6.85), Inches(legend_y - 0.02),
        Inches(1), Inches(0.25)
    )
    kd_label.text_frame.paragraphs[0].text = "Key Driver"
    kd_label.text_frame.paragraphs[0].font.size = Pt(9)
    kd_label.text_frame.paragraphs[0].font.color.rgb = COLORS["text_dark"]

    # Risk legend
    risk_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8), Inches(legend_y),
        Inches(0.3), Inches(0.2)
    )
    risk_shape.fill.solid()
    risk_shape.fill.fore_color.rgb = COLORS["risk"]
    risk_shape.line.fill.background()

    risk_label = slide.shapes.add_textbox(
        Inches(8.35), Inches(legend_y - 0.02),
        Inches(0.8), Inches(0.25)
    )
    risk_label.text_frame.paragraphs[0].text = "Risk Area"
    risk_label.text_frame.paragraphs[0].font.size = Pt(9)
    risk_label.text_frame.paragraphs[0].font.color.rgb = COLORS["text_dark"]


def add_summary_slide(prs: Presentation, company_name: str, metrics: Dict[str, str]):
    """Add summary slide with key metrics."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Metrics Summary"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS["text_dark"]

    if not metrics:
        return

    # Metrics grid
    cols = 3
    col_width = 3.0
    row_height = 1.2
    start_x = 0.5
    start_y = 1.0

    for i, (key, value) in enumerate(metrics.items()):
        col = i % cols
        row = i // cols

        x = start_x + col * col_width
        y = start_y + row * row_height

        # Metric box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(2.7), Inches(1.0)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS["bg_light"]
        box.line.color.rgb = COLORS["border"]

        # Metric label
        label_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + 0.1),
            Inches(2.5), Inches(0.3)
        )
        label_box.text_frame.paragraphs[0].text = key
        label_box.text_frame.paragraphs[0].font.size = Pt(11)
        label_box.text_frame.paragraphs[0].font.color.rgb = COLORS["text_dark"]

        # Metric value
        value_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + 0.45),
            Inches(2.5), Inches(0.4)
        )
        value_box.text_frame.paragraphs[0].text = str(value)
        value_box.text_frame.paragraphs[0].font.size = Pt(20)
        value_box.text_frame.paragraphs[0].font.bold = True
        value_box.text_frame.paragraphs[0].font.color.rgb = COLORS["primary"]


# ============================================================================
# Main API
# ============================================================================

def create_value_driver_pptx(
    industry: str,
    company_name: str = "Target Company",
    metrics: Dict[str, str] = None,
    output_path: str = None,
    include_ebitda: bool = True,
    include_working_capital: bool = True,
    include_capex: bool = True,
) -> str:
    """
    Create a PowerPoint with value driver diagrams.

    Args:
        industry: Industry name (must match key in INDUSTRY_TREES)
        company_name: Name of the target company
        metrics: Dict of metric names to values (e.g., {"ARR": "$42M"})
        output_path: Optional output file path
        include_ebitda: Include EBITDA tree slide
        include_working_capital: Include Working Capital tree slide
        include_capex: Include Capex tree slide

    Returns:
        Path to the generated PowerPoint file
    """
    metrics = metrics or {}

    # Normalize industry name
    industry_key = None
    for key in INDUSTRY_TREES.keys():
        if key.lower() == industry.lower() or key.lower().replace(" / ", "/") == industry.lower().replace(" / ", "/"):
            industry_key = key
            break

    if not industry_key:
        # Try partial match
        for key in INDUSTRY_TREES.keys():
            if industry.lower() in key.lower() or key.lower() in industry.lower():
                industry_key = key
                break

    if not industry_key:
        available = ", ".join(INDUSTRY_TREES.keys())
        raise ValueError(f"Unknown industry: {industry}. Available: {available}")

    trees = INDUSTRY_TREES[industry_key]

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    add_title_slide(prs, company_name, industry_key)

    # Tree slides
    if include_ebitda and "ebitda" in trees:
        tree_title = "EBITDA" if "Banking" not in industry_key else "PPNR"
        if "Insurance" in industry_key:
            tree_title = "Underwriting Income"
        if "Asset Management" in industry_key:
            tree_title = "Revenue & Earnings"
        add_tree_slide(prs, f"{tree_title} Value Driver Tree", trees["ebitda"], metrics)

    if include_working_capital and "working_capital" in trees:
        wc_title = "Working Capital Value Driver Tree"
        if "Banking" in industry_key:
            wc_title = "Balance Sheet Drivers"
        if "Insurance" in industry_key:
            wc_title = "Float & Reserves"
        add_tree_slide(prs, wc_title, trees["working_capital"], metrics)

    if include_capex and "capex" in trees:
        add_tree_slide(prs, "Capex Value Driver Tree", trees["capex"], metrics)

    # Summary slide with metrics
    if metrics:
        add_summary_slide(prs, company_name, metrics)

    # Save
    if not output_path:
        safe_name = re.sub(r'[^\w\s-]', '', company_name).replace(' ', '_')
        output_path = f"{safe_name}_Value_Drivers.pptx"

    prs.save(output_path)
    return output_path


def list_available_industries() -> List[str]:
    """Return list of available industries."""
    return list(INDUSTRY_TREES.keys())


# ============================================================================
# CLI / Testing
# ============================================================================

if __name__ == "__main__":
    # Example usage
    print("Available industries:", list_available_industries())

    # Generate sample for SaaS
    path = create_value_driver_pptx(
        industry="Technology / SaaS",
        company_name="Acme Software",
        metrics={
            "ARR": "$42M",
            "Growth": "+28% YoY",
            "Gross Margin": "78%",
            "NRR": "115%",
            "Customers": "2,500",
        }
    )
    print(f"Generated: {path}")
