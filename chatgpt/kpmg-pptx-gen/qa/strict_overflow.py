#!/usr/bin/env python3
"""Strict overflow check for PPTX outputs.

Renders an enlarged PPTX with padding and inspects the margins to detect
content overflow beyond the original canvas.
"""

from __future__ import annotations

import argparse
import json
import tempfile
from os.path import abspath, expanduser, join
from typing import Sequence, cast
from zipfile import ZipFile
import xml.etree.ElementTree as ET
from pathlib import Path
import sys

import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Emu

# Ensure repo root is on sys.path for qa imports when run as a script.
sys.path.append(str(Path(__file__).resolve().parents[1]))
from qa import render

PAD_PX: int = 100
PAD_RGB = (200, 200, 200)
EMU_PER_INCH: int = 914_400


def _calc_dpi_via_ooxml(input_path: str, max_w_px: int, max_h_px: int) -> int:
    """Calculate DPI from OOXML slide size to fit within max pixel bounds."""
    with ZipFile(input_path, "r") as zf:
        xml = zf.read("ppt/presentation.xml")
    root = ET.fromstring(xml)
    ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    sld_sz = root.find("p:sldSz", ns)
    if sld_sz is None:
        raise RuntimeError("Slide size not found in presentation.xml")
    cx = int(sld_sz.get("cx") or 0)
    cy = int(sld_sz.get("cy") or 0)
    if cx <= 0 or cy <= 0:
        raise RuntimeError("Invalid slide size values in presentation.xml")
    width_in = cx / EMU_PER_INCH
    height_in = cy / EMU_PER_INCH
    return round(min(max_w_px / width_in, max_h_px / height_in))


def _px_to_emu(px: int, dpi: int) -> Emu:
    """Convert pixels to EMU for a given DPI."""
    return Emu(int(px * EMU_PER_INCH // dpi))


def _calc_tol(dpi: int) -> int:
    """Calculate per-channel colour tolerance appropriate for DPI."""
    if dpi >= 300:
        return 0
    tol = round((300 - dpi) / 25)
    return min(max(tol, 1), 10)


def _enlarge_deck(src: str, dst: str, pad_emu: Emu) -> tuple[int, int]:
    """Enlarge the input PPTX with padding and return the new page size."""
    prs = Presentation(src)
    w0 = cast(Emu, prs.slide_width)
    h0 = cast(Emu, prs.slide_height)
    w1 = Emu(w0 + 2 * pad_emu)
    h1 = Emu(h0 + 2 * pad_emu)
    prs.slide_width = w1
    prs.slide_height = h1

    for slide in prs.slides:
        # Shift all shapes so the original canvas sits centered in the new deck.
        for shp in list(slide.shapes):
            shp.left = Emu(int(shp.left) + pad_emu)
            shp.top = Emu(int(shp.top) + pad_emu)

        pads = (
            (Emu(0), Emu(0), pad_emu, h1),
            (Emu(int(w1) - int(pad_emu)), Emu(0), pad_emu, h1),
            (Emu(0), Emu(0), w1, pad_emu),
            (Emu(0), Emu(int(h1) - int(pad_emu)), w1, pad_emu),
        )

        sp_tree = slide.shapes._spTree  # pylint: disable=protected-access

        for left, top, width, height in pads:
            pad_shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
            )
            pad_shape.fill.solid()
            pad_shape.fill.fore_color.rgb = RGBColor(*PAD_RGB)
            pad_shape.line.fill.background()

            # Send pad behind all other shapes (index 2 after mandatory nodes).
            sp_tree.remove(pad_shape._element)
            sp_tree.insert(2, pad_shape._element)

    prs.save(dst)
    return int(w1), int(h1)


def _inspect_images(
    paths: Sequence[str],
    pad_ratio_w: float,
    pad_ratio_h: float,
    dpi: int,
) -> list[int]:
    """Return 1-based indices of slides that contain pixels outside the pad."""

    tol = _calc_tol(dpi)
    failures: list[int] = []
    pad_colour = np.array(PAD_RGB, dtype=np.uint8)

    for idx, img_path in enumerate(paths, start=1):
        with Image.open(img_path) as img:
            rgb = img.convert("RGB")
            arr = np.asarray(rgb)

        h, w, _ = arr.shape
        pad_x = int(w * pad_ratio_w) - 1
        pad_y = int(h * pad_ratio_h) - 1

        left_margin = arr[:, :pad_x, :]
        right_margin = arr[:, w - pad_x :, :]
        top_margin = arr[:pad_y, :, :]
        bottom_margin = arr[h - pad_y :, :, :]

        def _is_clean(margin: np.ndarray) -> bool:
            diff = np.abs(margin.astype(np.int16) - pad_colour)
            matches = np.all(diff <= tol, axis=-1)
            mismatch_fraction = 1.0 - (np.count_nonzero(matches) / matches.size)
            if dpi >= 300:
                max_mismatch = 0.01
            elif dpi >= 200:
                max_mismatch = 0.02
            else:
                max_mismatch = 0.03
            return mismatch_fraction <= max_mismatch

        if not (
            _is_clean(left_margin)
            and _is_clean(right_margin)
            and _is_clean(top_margin)
            and _is_clean(bottom_margin)
        ):
            failures.append(idx)

    return failures


def main() -> None:
    parser = argparse.ArgumentParser(
        description=(
            "Check a PPTX for content overflowing the original canvas by rendering with padding "
            "and inspecting the margins."
        )
    )
    parser.add_argument("input_path", type=str, help="Path to the input PPTX file.")
    parser.add_argument("--out", required=True, help="Directory to write rendered PNGs and report.")
    parser.add_argument(
        "--width",
        type=int,
        default=1600,
        help=(
            "Approximate maximum width in pixels after isotropic scaling (default 1600). "
            "The actual value may exceed slightly."
        ),
    )
    parser.add_argument(
        "--height",
        type=int,
        default=900,
        help=(
            "Approximate maximum height in pixels after isotropic scaling (default 900). "
            "The actual value may exceed slightly."
        ),
    )
    parser.add_argument(
        "--pad_px",
        type=int,
        default=PAD_PX,
        help="Padding in pixels to add on each side before rasterization.",
    )
    args = parser.parse_args()

    input_path = abspath(expanduser(args.input_path))
    out_dir = Path(abspath(expanduser(args.out)))
    dpi = _calc_dpi_via_ooxml(input_path, args.width, args.height)

    with tempfile.TemporaryDirectory(prefix="strict_overflow_") as tmpdir:
        enlarged_pptx = join(tmpdir, "enlarged.pptx")
        pad_emu = _px_to_emu(args.pad_px, dpi)
        w1, h1 = _enlarge_deck(input_path, enlarged_pptx, pad_emu=pad_emu)
        pad_ratio_w = pad_emu / w1
        pad_ratio_h = pad_emu / h1

        img_paths = render.render_pptx_to_pngs_via_pdf(
            Path(enlarged_pptx), out_dir, dpi=dpi, prefix="slide"
        )
        failing = _inspect_images([str(p) for p in img_paths], pad_ratio_w, pad_ratio_h, dpi)

    report_path = str(out_dir / "overflow-report.json")
    report = {
        "input": input_path,
        "dpi": dpi,
        "failures": failing,
        "image_dir": str(out_dir),
    }
    with open(report_path, "w", encoding="utf-8") as f:
        json.dump(report, f, indent=2)

    if failing:
        print(
            "ERROR: Slides with content overflowing original canvas (1-based indexing): "
            + ", ".join(map(str, failing))
        )
        print(f"Report: {report_path}")
        raise SystemExit(1)

    print("Strict overflow check passed. No overflow detected.")


if __name__ == "__main__":
    main()
