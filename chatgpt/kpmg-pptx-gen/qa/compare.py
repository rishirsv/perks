from __future__ import annotations

import csv
import difflib
from pathlib import Path
from typing import Dict, Iterable, List, Optional, Tuple


def _require_optional_deps() -> None:
    try:
        import PIL  # noqa: F401
        import numpy  # noqa: F401
    except Exception as e:  # pragma: no cover
        raise RuntimeError(
            "Missing optional QA deps. Create a venv and install `requirements/requirements-qa.txt`.\n"
            "Example:\n"
            "  python3 -m venv .venv-qa\n"
            "  source .venv-qa/bin/activate\n"
            "  python -m pip install -r requirements/requirements-qa.txt"
        ) from e


def _load_png(path: Path):
    from PIL import Image

    with Image.open(path) as im:
        return im.convert("RGB")


def _diff_percent(expected_png: Path, actual_png: Path) -> float:
    import numpy as np

    exp = _load_png(expected_png)
    act = _load_png(actual_png)
    if exp.size != act.size:
        act = act.resize(exp.size)

    a = np.asarray(exp, dtype=np.int16)
    b = np.asarray(act, dtype=np.int16)
    d = np.abs(a - b).max(axis=2)
    return float((d > 10).mean() * 100.0)


def _write_compare_png(expected_png: Path, actual_png: Optional[Path], out_path: Path) -> None:
    from PIL import Image, ImageChops, ImageEnhance

    exp = _load_png(expected_png)
    if actual_png and actual_png.exists():
        act = _load_png(actual_png)
        if exp.size != act.size:
            act = act.resize(exp.size)
        diff = ImageChops.difference(exp, act)
        diff = ImageEnhance.Contrast(diff).enhance(4.0)
    else:
        act = Image.new("RGB", exp.size, "#FFFFFF")
        diff = Image.new("RGB", exp.size, "#FFFFFF")

    w, h = exp.size
    canvas = Image.new("RGB", (w * 3, h), "#FFFFFF")
    canvas.paste(exp, (0, 0))
    canvas.paste(act, (w, 0))
    canvas.paste(diff, (w * 2, 0))
    out_path.parent.mkdir(parents=True, exist_ok=True)
    canvas.save(out_path)


def _extract_titles_from_pptx(pptx_path: Path) -> List[str]:
    _require_optional_deps()
    from pptx import Presentation

    pres = Presentation(str(pptx_path))
    titles = []
    for slide in pres.slides:
        best_text = ""
        best_score = -1
        for sh in slide.shapes:
            if not getattr(sh, "has_text_frame", False):
                continue
            raw = ""
            try:
                raw = sh.text_frame.text or ""
            except Exception:
                raw = ""
            txt = " ".join(raw.split())
            if not txt:
                continue

            # Heuristic:
            # - Prefer larger text
            # - Prefer shapes nearer the top of the slide
            y_emu = int(getattr(sh, "top", 0) or 0)
            max_pt = 0
            try:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.size:
                            max_pt = max(max_pt, int(r.font.size.pt))
            except Exception:
                max_pt = 0

            # Higher score is better. `y_emu` is smaller when higher on the slide.
            score = (max_pt * 1_000_000) - y_emu
            if score > best_score:
                best_score = score
                best_text = txt

        titles.append(best_text)
    return titles


def _match_titles(expected: List[str], actual: List[str]) -> List[Optional[int]]:
    """
    For each expected title index, return the matched actual index (0-based), or None.
    Greedy matching by best similarity, enforcing 1:1 matches.
    """
    used = set()
    mapping: List[Optional[int]] = []
    for exp in expected:
        best = None
        best_score = 0.0
        for j, act in enumerate(actual):
            if j in used:
                continue
            score = difflib.SequenceMatcher(a=exp.lower(), b=act.lower()).ratio()
            if score > best_score:
                best_score = score
                best = j
        if best is not None and best_score >= 0.80:
            used.add(best)
            mapping.append(best)
        else:
            mapping.append(None)
    return mapping


def compare_png_dirs(
    expected_dir: Path,
    actual_dir: Path,
    out_dir: Path,
    *,
    expected_titles: Optional[List[str]] = None,
    actual_titles: Optional[List[str]] = None,
    match: str = "index",
) -> Path:
    """
    Compare two directories of per-slide PNGs and write:
    - `report.csv` with per-slide diff percentages
    - 3-up compare PNGs (expected | actual | diff)

    `match`:
    - "index": slide-01 vs slide-01, etc.
    - "title": requires `expected_titles` and `actual_titles` and matches by fuzzy title.
    """
    _require_optional_deps()

    expected_dir = Path(expected_dir)
    actual_dir = Path(actual_dir)
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    exp_pngs = sorted(expected_dir.glob("slide-*.png"))
    act_pngs = sorted(actual_dir.glob("slide-*.png"))

    def slide_num(p: Path) -> int:
        return int(p.stem.split("-")[-1])

    exp_by_num = {slide_num(p): p for p in exp_pngs}
    act_by_num = {slide_num(p): p for p in act_pngs}

    if match == "title":
        if expected_titles is None or actual_titles is None:
            raise ValueError("match='title' requires expected_titles and actual_titles")
        mapping = _match_titles(expected_titles, actual_titles)
    else:
        mapping = None

    report_path = out_dir / "report.csv"
    with report_path.open("w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        w.writerow(
            [
                "expected_slide",
                "expected_title",
                "actual_slide",
                "actual_title",
                "title_match_ratio",
                "diff_percent",
                "compare_png",
            ]
        )

        for exp_idx, exp_title in enumerate(expected_titles or [""] * len(exp_by_num), start=1):
            exp_png = exp_by_num.get(exp_idx)
            if not exp_png:
                continue

            if match == "title" and mapping is not None:
                act_idx0 = mapping[exp_idx - 1] if exp_idx - 1 < len(mapping) else None
                act_idx = (act_idx0 + 1) if act_idx0 is not None else None
            else:
                act_idx = exp_idx if exp_idx in act_by_num else None

            act_png = act_by_num.get(act_idx) if act_idx else None
            act_title = ""
            title_ratio = ""
            if act_idx and actual_titles and act_idx - 1 < len(actual_titles):
                act_title = actual_titles[act_idx - 1]
                if exp_title:
                    title_ratio = f"{difflib.SequenceMatcher(a=exp_title.lower(), b=act_title.lower()).ratio():.4f}"

            out_png = out_dir / (
                f"compare-exp{exp_idx:02d}-act{act_idx:02d}.png" if act_idx else f"compare-exp{exp_idx:02d}-missing.png"
            )
            _write_compare_png(exp_png, act_png, out_png)

            diff_pct = _diff_percent(exp_png, act_png) if act_png else 100.0
            w.writerow(
                [
                    exp_idx,
                    exp_title,
                    act_idx or "",
                    act_title,
                    title_ratio,
                    f"{diff_pct:.4f}",
                    str(out_png),
                ]
            )

    return report_path


def compare_decks_by_title(
    expected_pptx: Path,
    actual_pptx: Path,
    expected_png_dir: Path,
    actual_png_dir: Path,
    out_dir: Path,
) -> Path:
    exp_titles = _extract_titles_from_pptx(Path(expected_pptx))
    act_titles = _extract_titles_from_pptx(Path(actual_pptx))
    return compare_png_dirs(
        Path(expected_png_dir),
        Path(actual_png_dir),
        Path(out_dir),
        expected_titles=exp_titles,
        actual_titles=act_titles,
        match="title",
    )
