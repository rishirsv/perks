"""
Simple template analyzer for python-pptx files.

Usage:
    python tools/analyze_template.py glossary_template.pptx

Outputs basic shape info (index, type, position, size, and sample text) to
help define slide type mappings. Intended for one-time use when onboarding a
new template. Keep it lightweight and fast.

Example output:
    Slide 0 ------------------------------
    #00 type=TABLE name='Table 1' pos=(914400,914400) size=(3657600,5486400) text='Term | Definition'
    #01 type=TABLE name='Table 2' pos=(4699000,914400) size=(3657600,5486400) text='Term | Definition'
"""

import sys
from pathlib import Path
from pptx import Presentation


def describe_shape(idx, shape):
    kind = getattr(shape, "shape_type", "n/a")
    name = getattr(shape, "name", "")
    left = getattr(shape, "left", "?")
    top = getattr(shape, "top", "?")
    width = getattr(shape, "width", "?")
    height = getattr(shape, "height", "?")
    text = ""
    if getattr(shape, "has_text_frame", False):
        text = shape.text.strip().replace("\n", " | ")[:60]

    return f"#{idx:02d} type={kind} name='{name}' pos=({left},{top}) size=({width},{height}) text='{text}'"


def main(path):
    prs = Presentation(path)
    for s_idx, slide in enumerate(prs.slides):
        print(f"\nSlide {s_idx} ------------------------------")
        for idx, shape in enumerate(slide.shapes):
            print(describe_shape(idx, shape))


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python tools/analyze_template.py <template.pptx>")
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    if not pptx_path.exists():
        print(f"File not found: {pptx_path}")
        sys.exit(1)

    main(pptx_path)
