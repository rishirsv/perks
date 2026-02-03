"""Generate KPMG-branded glossary slides from JSON input."""

import argparse
import json
import sys
from pathlib import Path
from typing import Sequence, TypedDict, Union

# Add parent directory to path for imports
sys.path.insert(0, str(Path(__file__).parent.parent))

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

from core.rearrange import duplicate_slide
from core.formatting import (
    apply_text_formatting,
    apply_paragraph_formatting,
    clear_bullets,
)
from core.pagination import paginate_glossary_items

TERM_FORMAT = {"font_name": "Arial", "font_size": 10, "bold": True, "color": "00338D"}
DEFINITION_FORMAT = {"font_name": "Arial", "font_size": 10, "bold": False, "color": "000000"}
PARAGRAPH_FORMAT = {"bullet": False, "alignment": "LEFT"}


def _extract_term_definition(item):
    """Extract term and definition from item, supporting dict or tuple format."""
    if isinstance(item, dict):
        return item.get("term", ""), item.get("definition", "")
    elif isinstance(item, (list, tuple)) and len(item) >= 2:
        return str(item[0]), str(item[1])
    return "", ""


class GlossaryItem(TypedDict):
    term: str
    definition: str


def _find_tables(slide):
    return [shape.table for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.TABLE]


def _validate_tables(slide):
    tables = _find_tables(slide)
    if len(tables) < 2:
        raise ValueError(f"Template must contain at least 2 tables; found {len(tables)}.")

    for idx, table in enumerate(tables[:2], start=1):
        if len(table.columns) < 2:
            raise ValueError(f"Table {idx} must have at least 2 columns.")

    return tables


def _format_cell(cell, text, text_cfg):
    cell.text = text
    for paragraph in cell.text_frame.paragraphs:
        clear_bullets(paragraph)
        apply_paragraph_formatting(paragraph, PARAGRAPH_FORMAT)
        for run in paragraph.runs:
            apply_text_formatting(run, text_cfg)


def _clear_rows(table, start_row):
    for row_idx in range(start_row, len(table.rows)):
        for col in range(len(table.columns)):
            _format_cell(table.cell(row_idx, col), "", DEFINITION_FORMAT)


def _load_items(input_data: str):
    if input_data.endswith(".json") and Path(input_data).exists():
        with open(input_data, "r") as f:
            items = json.load(f)
    else:
        try:
            items = json.loads(input_data)
        except json.JSONDecodeError as exc:
            raise ValueError("Input must be a valid JSON string or path to JSON file.") from exc

    if not isinstance(items, list):
        raise ValueError("Input must be a list of dictionaries.")

    for idx, item in enumerate(items):
        if not isinstance(item, dict):
            raise ValueError(f"Item {idx} is not a dictionary.")

    return items


def _fill_slide(slide, page_data):
    """Fill a slide with glossary items using adaptive row heights.

    Args:
        slide: The slide to fill.
        page_data: Dict with 'left' and 'right' keys, each containing
            a list of (item, height_pt) tuples from pagination.
    """
    tables = _validate_tables(slide)

    left_table, right_table = tables[0], tables[1]
    left_items = page_data.get('left', [])
    right_items = page_data.get('right', [])

    # Fill left table with adaptive row heights
    for idx, (item, height_pt) in enumerate(left_items, start=1):
        term, definition = _extract_term_definition(item)
        left_table.rows[idx].height = Pt(height_pt)
        _format_cell(left_table.cell(idx, 0), term, TERM_FORMAT)
        _format_cell(left_table.cell(idx, 1), definition, DEFINITION_FORMAT)
    _clear_rows(left_table, len(left_items) + 1)

    # Fill right table with adaptive row heights
    for idx, (item, height_pt) in enumerate(right_items, start=1):
        term, definition = _extract_term_definition(item)
        right_table.rows[idx].height = Pt(height_pt)
        _format_cell(right_table.cell(idx, 0), term, TERM_FORMAT)
        _format_cell(right_table.cell(idx, 1), definition, DEFINITION_FORMAT)
    _clear_rows(right_table, len(right_items) + 1)


def generate_glossary(
    template_path: Union[Path, str],
    items: Sequence[GlossaryItem],
    output_path: Union[Path, str],
) -> Path:
    """Fill the glossary template; duplicate slides when capacity is exceeded.

    Uses dynamic pagination based on estimated text height to avoid overflow.

    Args:
        template_path: Path to the glossary template PPTX file.
        items: Sequence of glossary items. Each item can be:
            - dict with 'term' and 'definition' keys
            - tuple/list of (term, definition)
        output_path: Path where the generated PPTX will be saved.

    Returns:
        Path to the generated output file.

    Note:
        If items is empty, creates an output file with a single empty glossary
        slide (the template slide with cleared tables).
    """
    template_path = Path(template_path)
    output_path = Path(output_path)

    prs = Presentation(str(template_path))
    base_index = 0

    first_slide = prs.slides[base_index]
    tables = _validate_tables(first_slide)

    left_cap = len(tables[0].rows) - 1
    right_cap = len(tables[1].rows) - 1
    if left_cap <= 0 or right_cap <= 0:
        raise ValueError("Template tables have no available rows for glossary entries.")

    # Handle empty input explicitly: create one empty slide
    if not items:
        _fill_slide(first_slide, {'left': [], 'right': []})
        prs.save(str(output_path))
        print(f"Generated empty glossary (0 terms, 1 slide) at {output_path}")
        return output_path

    # Use height-based pagination for adaptive row heights
    pages = paginate_glossary_items(list(items))

    print(f"Paginated {len(items)} terms into {len(pages)} slide(s)")

    for page_idx, page_data in enumerate(pages):
        target_slide = first_slide if page_idx == 0 else duplicate_slide(prs, base_index)
        _fill_slide(target_slide, page_data)

    prs.save(str(output_path))
    print(f"Generated glossary with {len(items)} terms across {len(pages)} slide(s) at {output_path}")
    return output_path


def main():
    parser = argparse.ArgumentParser(description="Generate KPMG Glossary Slides")
    parser.add_argument("--template", required=True, help="Path to template PPTX")
    parser.add_argument("--input", required=True, help="JSON string or file path with terms")
    parser.add_argument("--output", required=True, help="Output PPTX path")

    args = parser.parse_args()

    try:
        items = _load_items(args.input)
        generate_glossary(args.template, items, args.output)
    except Exception as exc:  # pragma: no cover - surfaced to user
        print(str(exc), file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
