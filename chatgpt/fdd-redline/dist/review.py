from __future__ import annotations

import csv
import hashlib
import io
import json
import math
import re
from collections import Counter, defaultdict
from dataclasses import dataclass
from decimal import Decimal
from pathlib import Path
from typing import Any

from pptx import Presentation

from checks_crossref import run_crossref_checks
from checks_grammar import find_double_spaces, find_misspellings, find_repeated_words
from checks_placeholders import find_placeholders
from checks_terminology import find_term_variants
from extraction import extract_metrics_from_text, group_contiguous
from montage import create_montage as _create_montage
from models import DocumentModel, Finding, PageModel


@dataclass(slots=True)
class MetricValue:
    """A parsed metric value normalized for comparisons across pages."""

    kind: str  # "money" | "percent"
    normalized: Decimal


@dataclass(slots=True)
class SlideText:
    """Raw, unmodified text extracted from a single slide."""

    slide_num: int
    text: str


_CATEGORY_PREFIX: dict[str, str] = {
    "Placeholders": "P",
    "Grammar": "G",
    "Calculations": "C",
    "Style": "S",
    "Alignment": "A",
    "Cross-references": "X",
}

_GENERIC_METRIC_LABELS = {"consists", "includes", "comprises"}


def _normalize_comment_for_stable_id(comment: str) -> str:
    """Make comment text stable across minor whitespace/casing changes."""

    lower = comment.strip().lower()
    lower = re.sub(r"\s+", " ", lower)
    lower = re.sub(r"[^a-z0-9 ]+", "", lower)
    return lower.strip()


def _stable_id(page: int, category: str, comment: str) -> str:
    """Create a deterministic, short stable id like 'P-a3f2'."""

    prefix = _CATEGORY_PREFIX.get(category, "Z")
    payload = f"{page}|{category}|{_normalize_comment_for_stable_id(comment)}".encode("utf-8")
    digest = hashlib.sha1(payload).hexdigest()[:4]
    return f"{prefix}-{digest}"


def normalize_metric_value(raw: str) -> MetricValue:
    """Parse common money/percent formats into a comparable normalized value."""

    s = raw.strip()
    if not s:
        raise ValueError("Empty metric value")

    # Percent (convert 15.2% -> 0.152)
    percent_match = re.search(r"(-?\(?\s*)(\d[\d,]*\.?\d*)\s*%\s*\)?", s)
    if percent_match:
        sign_prefix = percent_match.group(1)
        number = Decimal(percent_match.group(2).replace(",", ""))
        is_negative = "-" in sign_prefix or "(" in sign_prefix
        if is_negative:
            number = -number
        return MetricValue(kind="percent", normalized=(number / Decimal("100")))

    # Money / numeric value with optional scale suffix.
    is_negative = False
    if "(" in s and ")" in s:
        is_negative = True

    cleaned = s
    cleaned = cleaned.replace("(", "").replace(")", "")
    cleaned = cleaned.strip()

    # Remove common currency codes/symbols.
    cleaned = re.sub(r"^(USD|CAD|EUR|GBP)\s+", "", cleaned, flags=re.IGNORECASE)
    cleaned = cleaned.lstrip("$£€").strip()

    # Detect magnitude suffixes/words.
    multiplier = Decimal("1")
    word_lower = cleaned.lower()
    if re.search(r"\b(billion|bn)\b", word_lower):
        multiplier = Decimal("1000000000")
        cleaned = re.sub(r"\b(billion|bn)\b", "", cleaned, flags=re.IGNORECASE).strip()
    elif re.search(r"\b(million)\b", word_lower):
        multiplier = Decimal("1000000")
        cleaned = re.sub(r"\b(million)\b", "", cleaned, flags=re.IGNORECASE).strip()
    elif re.search(r"\b(thousand)\b", word_lower):
        multiplier = Decimal("1000")
        cleaned = re.sub(r"\b(thousand)\b", "", cleaned, flags=re.IGNORECASE).strip()
    else:
        suffix_match = re.search(r"([kKmMbB]{1,2})\s*$", cleaned)
        if suffix_match:
            suffix = suffix_match.group(1).lower()
            if suffix in {"k"}:
                multiplier = Decimal("1000")
            elif suffix in {"m", "mm"}:
                multiplier = Decimal("1000000")
            elif suffix in {"b", "bn"}:
                multiplier = Decimal("1000000000")
            cleaned = re.sub(r"([kKmMbB]{1,2})\s*$", "", cleaned).strip()

    cleaned = cleaned.replace(",", "")
    number_match = re.search(r"-?\d[\d]*\.?\d*", cleaned)
    if not number_match:
        raise ValueError(f"Unparseable metric value: {raw}")

    number = Decimal(number_match.group(0))
    if is_negative:
        number = -abs(number)

    return MetricValue(kind="money", normalized=(number * multiplier))


def check_metric_consistency(doc: DocumentModel) -> list[Finding]:
    """Flag obvious metric drift across pages (e.g., Revenue differs on page 2 vs page 5)."""

    if doc.total_pages < 2:
        return []

    # label -> page_num -> list[raw_value]
    values_by_metric: dict[str, dict[int, list[str]]] = defaultdict(lambda: defaultdict(list))
    label_display: dict[str, str] = {}

    for page in doc.pages:
        for label, value in page.metrics:
            normalized_label = re.sub(r"\s+", " ", label.strip()).lower()
            if not normalized_label or normalized_label in _GENERIC_METRIC_LABELS:
                continue
            values_by_metric[normalized_label][page.page_num].append(value)
            label_display.setdefault(normalized_label, label.strip())

    findings: list[Finding] = []
    for metric_key, values_by_page in values_by_metric.items():
        if len(values_by_page) < 2:
            continue

        parsed_by_page: dict[int, MetricValue] = {}
        for page_num, raw_values in values_by_page.items():
            # Use the first value per page for now (keeps comparisons simple).
            raw_value = raw_values[0]
            try:
                parsed_by_page[page_num] = normalize_metric_value(raw_value)
            except ValueError:
                continue

        if len(parsed_by_page) < 2:
            continue

        # Compare normalized values; if any disagree, flag a single Calculations finding.
        items = sorted(parsed_by_page.items(), key=lambda kv: kv[0])
        baseline_page, baseline_val = items[0]
        disagreements: list[tuple[int, str]] = [(baseline_page, values_by_page[baseline_page][0])]
        for page_num, parsed in items[1:]:
            if parsed.kind != baseline_val.kind:
                continue
            if parsed.normalized != baseline_val.normalized:
                disagreements.append((page_num, values_by_page[page_num][0]))

        if len(disagreements) <= 1:
            continue

        metric_name = label_display.get(metric_key, metric_key)
        examples = ", ".join([f"p{p}={v}" for p, v in disagreements[:6]])
        findings.append(
            Finding(
                id="",
                page=0,
                category="Calculations",
                priority="H",
                comment=f'Metric "{metric_name}" appears inconsistent across pages ({examples}).',
            )
        )

    return findings


def select_pages_for_visual_review(doc: DocumentModel, findings: list[Finding]) -> list[int]:
    """
    Select a capped set of pages worth rendering to PNG for layout/alignment review.

    The cap is set to 50% of the deck by default to avoid expensive full renders.
    """

    cap = max(1, int(math.floor(doc.total_pages * 0.5)))

    pages_with_findings = {f.page for f in findings if f.page > 0 and f.priority in {"H", "M"}}
    candidates: set[int] = set()
    for page in doc.pages:
        word_count = len(re.findall(r"\b\w+\b", page.raw_text))
        if page.has_table or page.has_chart:
            candidates.add(page.page_num)
        elif word_count < 30:
            candidates.add(page.page_num)

    candidates |= pages_with_findings

    def score(page_num: int) -> tuple[int, int]:
        page = doc.pages[page_num - 1]
        s = 0
        if page.has_table or page.has_chart:
            s += 3
        if page_num in pages_with_findings:
            s += 2
        word_count = len(re.findall(r"\b\w+\b", page.raw_text))
        if word_count < 30:
            s += 1
        return (-s, page_num)

    selected = sorted(candidates, key=score)[:cap]
    return selected


def compress_placeholder_findings(
    findings: list[Finding],
    *,
    max_distinct_tokens: int,
    top_n: int,
    other_token_list_limit: int,
) -> list[Finding]:
    """
    Compress large numbers of placeholder findings into a smaller "top tokens" list.

    This keeps the issue log scannable while still highlighting the most common
    placeholder problems.
    """

    placeholder_findings = [f for f in findings if f.category == "Placeholders"]
    other_findings = [f for f in findings if f.category != "Placeholders"]

    token_re = re.compile(r"\"([^\"]+)\"")
    tokens: list[str] = []
    pages_by_token: dict[str, set[int]] = defaultdict(set)
    for f in placeholder_findings:
        m = token_re.search(f.comment)
        if not m:
            continue
        token = m.group(1)
        tokens.append(token)
        if f.page > 0:
            pages_by_token[token].add(f.page)

    distinct = sorted(set(tokens))
    if len(distinct) <= max_distinct_tokens:
        return findings

    counter = Counter(tokens)
    top_tokens = [t for t, _ in counter.most_common(top_n)]
    compressed: list[Finding] = []
    for t in top_tokens:
        pages = sorted(pages_by_token.get(t, set()))
        ranges = group_contiguous(pages)
        if ranges:
            page_str = ", ".join([f"{a}" if a == b else f"{a}-{b}" for a, b in ranges])
            comment = f'Placeholder "{t}" found on pages {page_str}.'
        else:
            comment = f'Placeholder "{t}" found.'
        compressed.append(
            Finding(
                id="",
                page=0,
                category="Placeholders",
                priority="H",
                comment=comment,
            )
        )

    remaining = [t for t in distinct if t not in set(top_tokens)]
    remaining_list = ", ".join(remaining[:other_token_list_limit])
    remaining_more = len(remaining) - min(len(remaining), other_token_list_limit)
    if remaining_more > 0:
        remaining_list = f"{remaining_list} (+{remaining_more} more)"

    compressed.append(
        Finding(
            id="",
            page=0,
            category="Placeholders",
            priority="H",
            comment=f"Additional placeholder tokens detected: {remaining_list}.",
        )
    )

    return other_findings + compressed


def validate_findings(findings: list[Finding], doc: DocumentModel) -> list[Finding]:
    """
    De-duplicate, mark invalid items, and assign sequential row ids.

    Invalid means: a page number outside the document range (except 0 = global).
    """

    seen: set[tuple[int, str, str, str]] = set()
    deduped: list[Finding] = []
    for f in findings:
        key = (f.page, f.category, f.priority, _normalize_comment_for_stable_id(f.comment))
        if key in seen:
            continue
        seen.add(key)
        deduped.append(f)

    valid: list[Finding] = []
    invalid: list[Finding] = []
    for f in deduped:
        if f.page == 0:
            f.is_valid = True
            valid.append(f)
            continue
        if 1 <= f.page <= doc.total_pages:
            f.is_valid = True
            valid.append(f)
        else:
            f.is_valid = False
            invalid.append(f)

    valid_sorted = sorted(valid, key=lambda f: (f.page, f.category, f.priority, f.comment))
    out = valid_sorted + invalid

    for idx, f in enumerate(out, start=1):
        f.id = f"{idx:03d}"
        f.stable_id = _stable_id(f.page, f.category, f.comment)

    return out


def _filter_ignored_lines(raw_text: str, ignore: list[str]) -> str:
    """Remove known slide-master artifact lines before running checks."""

    ignore_set = {s.strip().lower() for s in ignore if str(s).strip()}
    if not ignore_set:
        return raw_text

    kept: list[str] = []
    for line in raw_text.splitlines():
        if line.strip().lower() in ignore_set:
            continue
        kept.append(line)
    return "\n".join(kept)


def _load_rules_text(path: Path) -> dict[str, Any]:
    """
    Load rules from rules.yaml.

    The runtime avoids mandatory external deps; it supports:
    - PyYAML when installed
    - JSON (valid YAML) fallback
    - A minimal YAML subset (lists + flat key/value maps)
    """

    text = path.read_text(encoding="utf-8")

    try:
        import yaml  # type: ignore

        data = yaml.safe_load(text)
        return data if isinstance(data, dict) else {}
    except Exception:
        pass

    try:
        data = json.loads(text)
        return data if isinstance(data, dict) else {}
    except json.JSONDecodeError:
        pass

    # Minimal YAML subset parser (enough for our shipped rules.yaml).
    rules: dict[str, Any] = {}
    current_key: str | None = None
    for raw_line in text.splitlines():
        line = raw_line.rstrip()
        if not line or line.lstrip().startswith("#"):
            continue

        top_match = re.match(r"^([A-Za-z_][A-Za-z0-9_-]*)\s*:\s*$", line)
        if top_match:
            current_key = top_match.group(1)
            rules[current_key] = []
            continue

        if current_key is None:
            continue

        if current_key in {"ignore", "terminology"}:
            item_match = re.match(r"^\s*-\s*(.+?)\s*$", line)
            if not item_match:
                continue
            item_raw = item_match.group(1).strip()
            item = item_raw.strip("\"' ")
            if current_key == "terminology":
                try:
                    rules[current_key].append(json.loads(item_raw))
                except Exception:
                    rules[current_key].append(item)
            else:
                rules[current_key].append(item)
            continue

        if current_key == "placeholders":
            if not isinstance(rules[current_key], list):
                rules[current_key] = []

            # List item start
            item_match = re.match(r"^\s*-\s*pattern\s*:\s*(.+?)\s*$", line)
            if item_match:
                pattern_raw = item_match.group(1).strip()
                pattern = pattern_raw.strip("\"' ")
                rules[current_key].append({"pattern": pattern})
                continue

            kv_match = re.match(r"^\s*(type|canonical)\s*:\s*(.+?)\s*$", line)
            if kv_match and rules[current_key]:
                k = kv_match.group(1)
                v_raw = kv_match.group(2).strip()
                v = v_raw.strip("\"' ")
                rules[current_key][-1][k] = v
            continue

        # Flat maps (misspellings/style)
        if current_key in {"misspellings", "style"}:
            if not isinstance(rules.get(current_key), dict):
                rules[current_key] = {}
            kv_match = re.match(r"^\s*([^:]+?)\s*:\s*(.+?)\s*$", line)
            if not kv_match:
                continue
            k = kv_match.group(1).strip().strip("\"' ")
            v = kv_match.group(2).strip().strip("\"' ")
            rules[current_key][k] = v
            continue

    return rules


def _rules_to_terminology_groups(rules: dict[str, Any]) -> dict[str, dict[str, object]]:
    """Convert terminology config into the structure used by `find_term_variants`."""

    raw_groups = rules.get("terminology", [])
    if not isinstance(raw_groups, list):
        return {}

    groups: dict[str, dict[str, object]] = {}
    for group in raw_groups:
        if not isinstance(group, list):
            continue
        variants = [str(v) for v in group if str(v).strip()]
        if len(variants) < 2:
            continue
        key = re.sub(r"[^a-z0-9]+", "", variants[0].lower())
        full_form = variants[-1]
        groups[key] = {"full_form": full_form, "abbreviations": variants}
    return groups


def _extract_pptx(path: str) -> DocumentModel:
    """Extract raw text from a PPTX deck slide-by-slide."""

    prs = Presentation(path)
    pages: list[PageModel] = []

    for idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        has_table = False
        has_chart = False
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                has_table = True
                continue
            if getattr(shape, "has_chart", False):
                has_chart = True

            if not hasattr(shape, "text_frame") or shape.text_frame is None:
                continue
            text = shape.text_frame.text
            if not text:
                continue
            parts.append(text)

        raw_text = "\n".join(parts)
        metrics = extract_metrics_from_text(raw_text)
        pages.append(
            PageModel(
                page_num=idx,
                raw_text=raw_text,
                has_table=has_table,
                has_chart=has_chart,
                metrics=metrics,
            )
        )

    return DocumentModel(
        pages=pages,
        source_path=path,
        source_type="pptx",
        total_pages=len(pages),
    )


def extract_slides(pptx_path: str) -> list[SlideText]:
    """
    Extract slide text for a PPTX deck.

    This is a lightweight helper intended for GPT orchestration: it returns
    unmodified text per slide and skips tables (visual review is better there).
    """

    prs = Presentation(pptx_path)
    slides: list[SlideText] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                continue
            if not hasattr(shape, "text_frame") or shape.text_frame is None:
                continue
            text = shape.text_frame.text
            if text:
                parts.append(text)
        slides.append(SlideText(slide_num=idx, text="\n".join(parts)))
    return slides


def run_checks(slides: list[SlideText], rules: dict[str, Any]) -> list[Finding]:
    """
    Run deterministic checks over extracted slide text and return findings.

    This mirrors the "comment-only" review behavior in the plan/spec.
    """

    pages: list[PageModel] = []
    for slide in slides:
        pages.append(
            PageModel(
                page_num=slide.slide_num,
                raw_text=slide.text,
                has_table=False,
                has_chart=False,
                metrics=extract_metrics_from_text(slide.text),
            )
        )
    doc = DocumentModel(
        pages=pages,
        source_path="(in-memory)",
        source_type="pptx",
        total_pages=len(pages),
    )

    ignore = rules.get("ignore", [])
    ignore_list = [str(s) for s in ignore] if isinstance(ignore, list) else []
    placeholder_patterns = rules.get("placeholders", [])
    patterns = placeholder_patterns if isinstance(placeholder_patterns, list) else []
    misspellings_raw = rules.get("misspellings", {})
    misspellings = misspellings_raw if isinstance(misspellings_raw, dict) else {}
    term_groups = _rules_to_terminology_groups(rules)

    findings: list[Finding] = []
    for page in doc.pages:
        page_for_checks = PageModel(
            page_num=page.page_num,
            raw_text=_filter_ignored_lines(page.raw_text, ignore_list),
            has_table=page.has_table,
            has_chart=page.has_chart,
            metrics=page.metrics,
        )
        findings.extend(find_placeholders(page_for_checks, patterns))
        findings.extend(find_repeated_words(page_for_checks))
        findings.extend(find_double_spaces(page_for_checks))
        findings.extend(find_misspellings(page_for_checks, misspellings))

    findings.extend(find_term_variants(doc, term_groups))
    findings.extend(run_crossref_checks(doc))
    findings.extend(check_metric_consistency(doc))

    return validate_findings(findings, doc)


def render_pngs(pdf_path: str) -> list[str]:
    """
    Render PDF pages to PNG files.

    This is optional and requires extra dependencies (e.g., `pdf2image` +
    Poppler). If unavailable, this raises a clear error.
    """

    try:
        from pdf2image import convert_from_path  # type: ignore
    except Exception as e:
        raise RuntimeError(
            "PDF rendering is unavailable. Install `pdf2image` and Poppler to enable it."
        ) from e

    images = convert_from_path(pdf_path)
    out_dir = Path(pdf_path).with_suffix("").with_name(f"{Path(pdf_path).stem}_pngs")
    out_dir.mkdir(parents=True, exist_ok=True)

    out_paths: list[str] = []
    for idx, img in enumerate(images, start=1):
        out_path = out_dir / f"page_{idx:03d}.png"
        img.save(out_path, format="PNG")
        out_paths.append(str(out_path))
    return out_paths


def create_montage(png_paths: list[str]) -> str:
    """
    Create a contact-sheet montage from PNG page images.

    Returns the montage image path on disk.
    """

    from PIL import Image

    if not png_paths:
        raise ValueError("No PNGs provided")

    out_path = Path(png_paths[0]).with_name("montage.png")
    cell_w = 360
    with Image.open(png_paths[0]) as first:
        ratio = first.height / max(1, first.width)
        cell_h = max(1, round(cell_w * ratio))

    _create_montage(
        input_files=png_paths,
        output_file=str(out_path),
        num_col=4,
        cell_w=cell_w,
        cell_h=cell_h,
        gap=14,
        label_mode="number",
    )
    return str(out_path)


def _issue_log_markdown(findings: list[Finding]) -> str:
    """Render findings as a markdown table (primary output)."""

    lines = [
        "| row_id | stable_id | Page | Category | Priority | Comment |",
        "|--------|-----------|------|----------|----------|---------|",
    ]
    for f in findings:
        page_str = "" if f.page == 0 else str(f.page)
        safe_comment = f.comment.replace("\n", " ").strip()
        lines.append(f"| {f.id} | {f.stable_id} | {page_str} | {f.category} | {f.priority} | {safe_comment} |")
    return "\n".join(lines) + "\n"


def _issue_log_csv(findings: list[Finding]) -> str:
    """Render findings as CSV (primary output)."""

    buf = io.StringIO()
    writer = csv.writer(buf)
    writer.writerow(["row_id", "stable_id", "Page", "Category", "Priority", "Comment"])
    for f in findings:
        page_str = "" if f.page == 0 else str(f.page)
        writer.writerow([f.id, f.stable_id, page_str, f.category, f.priority, f.comment])
    return buf.getvalue()


def run_review(source_path: str, *, render_visuals: bool = True) -> dict[str, Any]:
    """
    Run a full review of an FDD deck (PPTX today; PDF support is optional).

    Returns a dict with:
    - doc: DocumentModel
    - findings: list[Finding]
    - issue_log_md / issue_log_csv: formatted outputs
    - selected_pages_for_visual_review: list[int]
    """

    src = Path(source_path)
    ext = src.suffix.lower().lstrip(".")

    rules_path = Path(__file__).with_name("rules.yaml")
    rules = _load_rules_text(rules_path) if rules_path.exists() else {}

    if ext == "pptx":
        doc = _extract_pptx(str(src))
    elif ext == "pdf":
        raise RuntimeError("PDF input is not supported in this environment yet.")
    else:
        raise ValueError(f"Unsupported input type: {src.suffix}")

    ignore = rules.get("ignore", [])
    ignore_list = [str(s) for s in ignore] if isinstance(ignore, list) else []
    placeholder_patterns = rules.get("placeholders", [])
    patterns = placeholder_patterns if isinstance(placeholder_patterns, list) else []
    misspellings_raw = rules.get("misspellings", {})
    misspellings = misspellings_raw if isinstance(misspellings_raw, dict) else {}
    term_groups = _rules_to_terminology_groups(rules)

    findings: list[Finding] = []

    # Page-level checks.
    for page in doc.pages:
        page_for_checks = PageModel(
            page_num=page.page_num,
            raw_text=_filter_ignored_lines(page.raw_text, ignore_list),
            has_table=page.has_table,
            has_chart=page.has_chart,
            metrics=page.metrics,
        )
        findings.extend(find_placeholders(page_for_checks, patterns))
        findings.extend(find_repeated_words(page_for_checks))
        findings.extend(find_double_spaces(page_for_checks))
        findings.extend(find_misspellings(page_for_checks, misspellings))

    # Document-level checks.
    findings.extend(find_term_variants(doc, term_groups))
    findings.extend(run_crossref_checks(doc))
    findings.extend(check_metric_consistency(doc))

    # Optional compression to keep the output readable on large decks.
    findings = compress_placeholder_findings(
        findings,
        max_distinct_tokens=250,
        top_n=25,
        other_token_list_limit=50,
    )

    findings = validate_findings(findings, doc)
    selected = select_pages_for_visual_review(doc, findings)

    result: dict[str, Any] = {
        "doc": doc,
        "findings": findings,
        "issue_log_md": _issue_log_markdown(findings),
        "issue_log_csv": _issue_log_csv(findings),
        "selected_pages_for_visual_review": selected,
    }

    # Visual rendering is a nice-to-have; it is disabled in tests and in minimal environments.
    if render_visuals:
        result["visuals"] = None

    return result
